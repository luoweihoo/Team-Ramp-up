# @File    :   create_summary.py
# @Time    :   2019/09/05 15:48:28
# @Author  :   Wei Luo 
# @Version :   1.0
# @Contact :   luoweihoo@yahoo.com
# @Desc    :   Create the summary into a PPT slides deck

import openpyxl, sys, trfunction
from openpyxl.styles import Font
from openpyxl.styles.colors import BLACK
from openpyxl.styles.colors import RED
from pptx import Presentation
from pptx.util import Inches


trackList = 'Team Ramp-up Tracking.xlsx'
# Check if the track list is in the current folder; if yes open it
wb = trfunction.checkTrackList(trackList)
# Build the summary list
summaryList = trfunction.buildSummaryList(wb)

# Writ the result into a Excel List
wbSummary = 'SummaryList.xlsx'
listWb = openpyxl.Workbook()
listSheet = listWb.active

listSheet['A1'] = 'Manager'
listSheet['B1']  = '# of Nominated'
listSheet['C1']  = '# Quiz Taken'
listSheet['D1']  = '# Passed'
listSheet['E1']  = '# Failed'

rowNum = 2
for listRow in summaryList:
    listSheet['A' + str(rowNum)] = listRow[0]
    listSheet['B' + str(rowNum)] = listRow[1]
    listSheet['C' + str(rowNum)] = listRow[2]
    listSheet['D' + str(rowNum)] = listRow[3]
    listSheet['E' + str(rowNum)] = listRow[4]
    rowNum += 1

listWb.save(wbSummary)
print('Processing is finished!')

# prs = Presentation()
# title_only_slide_layout = prs.slide_layouts[5]
# slide = prs.slides.add_slide(title_only_slide_layout)
# shapes = slide.shapes
# shapes.title.text = 'Workshop List'

# rows = 4; cols = 3
# left = top = Inches(2.0)
# width = Inches(6.0)
# height = Inches(1.2)

# table = shapes.add_table(rows, cols, left, top, width, height).table
# table.columns[0].width = Inches(2.0)
# table.columns[1].width = Inches(2.0)
# table.columns[2].width = Inches(2.0)

# table.cell(0, 0).text = 'Sr. No.'
# table.cell(0, 1).text = 'Student Name'
# table.cell(0, 2).text = 'Score'

# lineNum = 0
# for testRow in (testResult[0:2]):
#     table.cell(lineNum + 1, 0).text = testRow[0]
#     table.cell(lineNum + 1, 1).text = testRow[1]
#     table.cell(lineNum + 1, 2).text = testRow[2]
#     lineNum += 1

# prs.save('TestResult.pptx')
# print('Processing is finished!')
